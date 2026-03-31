"""
Create a PowerPoint report from Flow360 post-processing figures.

Scans <figures_dir> for the three standard subfolders (total_force, residual,
forcehistory) and builds a .pptx with a title slide followed by a section
header + one slide per figure for each subfolder.

Usage:
    python create_ppt_report.py <case_name> [output.pptx]

    python create_ppt_report.py XV15_MRF
    python create_ppt_report.py XV15_MRF my_report.pptx

Figures are read from release_test/figures/<case_name>/.
Output is saved to release_test/report/<case_name>.pptx unless overridden.

Install dependency if needed:
    pip install python-pptx
"""

import datetime
import os
import re
import sys

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN


def _natural_key(filename):
    """Sort key that orders embedded integers numerically (e.g. RPM800 < RPM1200)."""
    return [int(c) if c.isdigit() else c.lower()
            for c in re.split(r'(\d+)', filename)]


# ── slide dimensions (widescreen 16:9) ────────────────────────────────────────
SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

# ── section metadata ──────────────────────────────────────────────────────────
SECTIONS = [
    ("user_define",  "User Define"),
    ("total_force",  "Total Force Coefficients"),
    ("runtime",      "Runtime Summary"),
    ("residual",     "Residuals"),
    ("forcehistory", "Force History"),
]

# ── colour palette ────────────────────────────────────────────────────────────
COLOR_TITLE_BG   = RGBColor(0x1F, 0x49, 0x7D)   # dark blue
COLOR_SECTION_BG = RGBColor(0x2E, 0x75, 0xB6)   # medium blue
COLOR_WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_DARK       = RGBColor(0x1F, 0x1F, 0x1F)


def _blank_slide(prs):
    """Add a completely blank slide and return it."""
    blank_layout = prs.slide_layouts[6]  # layout index 6 = blank
    return prs.slides.add_slide(blank_layout)


def _fill_bg(slide, color):
    """Fill the slide background with a solid colour."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_textbox(slide, text, left, top, width, height,
                 font_size=24, bold=False, color=COLOR_DARK,
                 align=PP_ALIGN.LEFT, wrap=True):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txBox


def add_title_slide(prs, case_name):
    slide = _blank_slide(prs)
    _fill_bg(slide, COLOR_TITLE_BG)

    # Main title
    _add_textbox(
        slide, f"Flow360 Post-Processing Report",
        left=Inches(1), top=Inches(2.2),
        width=Inches(11.33), height=Inches(1.2),
        font_size=40, bold=True, color=COLOR_WHITE, align=PP_ALIGN.CENTER,
    )
    # Sub-title: case name
    _add_textbox(
        slide, case_name,
        left=Inches(1), top=Inches(3.5),
        width=Inches(11.33), height=Inches(0.8),
        font_size=28, bold=False, color=COLOR_WHITE, align=PP_ALIGN.CENTER,
    )
    # Section list
    section_names = "  |  ".join(title for _, title in SECTIONS)
    _add_textbox(
        slide, section_names,
        left=Inches(1), top=Inches(4.6),
        width=Inches(11.33), height=Inches(0.6),
        font_size=16, bold=False, color=RGBColor(0xBD, 0xD7, 0xEE),
        align=PP_ALIGN.CENTER,
    )
    # Date
    date_str = datetime.date.today().strftime("%B %d, %Y")
    _add_textbox(
        slide, date_str,
        left=Inches(1), top=Inches(5.4),
        width=Inches(11.33), height=Inches(0.5),
        font_size=16, bold=False, color=RGBColor(0xBD, 0xD7, 0xEE),
        align=PP_ALIGN.CENTER,
    )


def add_section_header(prs, section_title):
    slide = _blank_slide(prs)
    _fill_bg(slide, COLOR_SECTION_BG)
    _add_textbox(
        slide, section_title,
        left=Inches(1), top=Inches(3.0),
        width=Inches(11.33), height=Inches(1.2),
        font_size=36, bold=True, color=COLOR_WHITE, align=PP_ALIGN.CENTER,
    )


def add_figure_slide(prs, img_path):
    slide = _blank_slide(prs)

    # Place image to fill the full slide
    img_top    = Inches(0.1)
    img_height = SLIDE_H - Inches(0.2)
    img_width  = SLIDE_W - Inches(0.4)
    _place_picture(slide, img_path, Inches(0.2), img_top, img_width, img_height)


def _place_picture(slide, img_path, left, top, max_width, max_height):
    """Add a picture inside a bounding box, preserving the original aspect ratio, centred."""
    # Add without explicit dimensions to read the natural (unscaled) size
    pic = slide.shapes.add_picture(img_path, left, top)
    nat_w, nat_h = pic.width, pic.height
    if nat_w > 0 and nat_h > 0:
        ratio = nat_w / nat_h
        new_h = max_height
        new_w = int(new_h * ratio)
        if new_w > max_width:
            new_w = max_width
            new_h = int(new_w / ratio)
        pic.width  = new_w
        pic.height = new_h
        pic.left   = left + int((max_width  - new_w) / 2)
        pic.top    = top  + int((max_height - new_h) / 2)


def add_dual_figure_slide(prs, img_path_left, img_path_right):
    """Place two figures side-by-side on a single slide."""
    slide = _blank_slide(prs)

    padding    = Inches(0.15)
    top        = Inches(0.1)
    half_w     = (SLIDE_W - padding * 3) // 2
    img_height = SLIDE_H - top - Inches(0.1)

    for i, img_path in enumerate((img_path_left, img_path_right)):
        left = padding + i * (half_w + padding)
        _place_picture(slide, img_path, left, top, half_w, img_height)


def main():
    if len(sys.argv) < 2:
        print(__doc__)
        sys.exit(1)

    case_name   = sys.argv[1]
    figures_dir = os.path.join("release_test", "figures", case_name)

    if len(sys.argv) > 2:
        output_path = sys.argv[2]
    else:
        report_dir  = os.path.join("release_test", "report")
        os.makedirs(report_dir, exist_ok=True)
        output_path = os.path.join(report_dir, f"{case_name}.pptx")
    print(f"Case:   {case_name}")
    print(f"Input:  {figures_dir}")
    print(f"Output: {output_path}")


    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    add_title_slide(prs, case_name)

    total_figures = 0
    for subfolder, section_title in SECTIONS:
        section_dir = os.path.join(figures_dir, subfolder)
        if not os.path.isdir(section_dir):
            print(f"  [skip] {subfolder}/ not found")
            continue

        _sort_key = _natural_key if subfolder in ("total_force", "residual") else str
        images = sorted(
            (f for f in os.listdir(section_dir)
             if f.lower().endswith((".png", ".jpg", ".jpeg"))),
            key=_sort_key,
        )
        if not images:
            print(f"  [skip] {subfolder}/ has no images")
            continue

        print(f"\n  Section: {section_title}  ({len(images)} figures)")
        add_section_header(prs, section_title)

        if subfolder == "total_force":
            # Pair 'compare' and 'diff' figures onto one slide; show the rest normally
            compare_img = next((os.path.join(section_dir, f) for f in images if "compare" in f.lower()), None)
            diff_img    = next((os.path.join(section_dir, f) for f in images if "diff"    in f.lower()), None)
            paired = {os.path.basename(p) for p in (compare_img, diff_img) if p}

            for img_file in images:
                if img_file in paired:
                    continue          # handled separately below
                img_path = os.path.join(section_dir, img_file)
                print(f"    + {img_file}")
                add_figure_slide(prs, img_path)
                total_figures += 1

            if compare_img and diff_img:
                print(f"    + [paired] {os.path.basename(compare_img)}  |  {os.path.basename(diff_img)}")
                add_dual_figure_slide(prs, compare_img, diff_img)
                total_figures += 1
            else:
                for img_path in filter(None, (compare_img, diff_img)):
                    print(f"    + {os.path.basename(img_path)}")
                    add_figure_slide(prs, img_path)
                    total_figures += 1
        elif subfolder == "user_define":
            # Pair 'compare' and 'diff' figures onto one slide; show the rest normally
            compare_img = next((os.path.join(section_dir, f) for f in images if "compare" in f.lower()), None)
            diff_img    = next((os.path.join(section_dir, f) for f in images if "diff"    in f.lower()), None)
            paired = {os.path.basename(p) for p in (compare_img, diff_img) if p}

            for img_file in images:
                if img_file in paired:
                    continue
                img_path = os.path.join(section_dir, img_file)
                print(f"    + {img_file}")
                add_figure_slide(prs, img_path)
                total_figures += 1

            if compare_img and diff_img:
                print(f"    + [paired] {os.path.basename(compare_img)}  |  {os.path.basename(diff_img)}")
                add_dual_figure_slide(prs, compare_img, diff_img)
                total_figures += 1
            else:
                for img_path in filter(None, (compare_img, diff_img)):
                    print(f"    + {os.path.basename(img_path)}")
                    add_figure_slide(prs, img_path)
                    total_figures += 1
        elif subfolder in ("residual", "forcehistory"):
            # Split into non-range (left) and range (right), pair by sort order
            no_range = sorted((f for f in images if "range" not in f.lower()), key=_natural_key)
            range_   = sorted((f for f in images if "range"     in f.lower()), key=_natural_key)
            unpaired = sorted(f for f in images
                              if f not in no_range and f not in range_)

            for left_f, right_f in zip(no_range, range_):
                left_p  = os.path.join(section_dir, left_f)
                right_p = os.path.join(section_dir, right_f)
                print(f"    + [paired] {left_f}  |  {right_f}")
                add_dual_figure_slide(prs, left_p, right_p)
                total_figures += 1

            # Any leftovers (uneven counts) get their own slide
            for img_file in list(no_range[len(range_):]) + list(range_[len(no_range):]) + unpaired:
                img_path = os.path.join(section_dir, img_file)
                print(f"    + {img_file}")
                add_figure_slide(prs, img_path)
                total_figures += 1
        else:
            for img_file in images:
                img_path = os.path.join(section_dir, img_file)
                print(f"    + {img_file}")
                add_figure_slide(prs, img_path)
                total_figures += 1

    prs.save(output_path)
    print(f"\nSaved: {output_path}  ({total_figures} figure slides)")


if __name__ == "__main__":
    main()
